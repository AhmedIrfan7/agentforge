"""PPTX text/structure extraction (roadmap step 092) -- registered into
extraction.py:HANDLERS under "pptx".

python-pptx (MIT). A slide's title shape becomes a markdown heading;
every other text-frame shape becomes a plain paragraph; tables become
markdown tables (extraction_tables.py, shared with pdf/docx/xlsx).
Title detection checks both PP_PLACEHOLDER.TITLE and CENTER_TITLE --
verified against a real generated pptx that a title-slide layout's
title shape is CENTER_TITLE, not TITLE, so checking only one of the two
would silently miss every title-slide-style title.

Shapes are read in slide.shapes order (insertion/z-order), not a
claimed visual reading order -- pptx has no single canonical reading
order for freely positioned shapes, and this doesn't pretend to solve
that, same "honest about what it actually does" stance as
extraction_pdf.py's font-size heading heuristic.
"""

import io

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.shapes.base import BaseShape

from extraction_tables import rows_to_markdown

_TITLE_PLACEHOLDER_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


def _is_title_shape(shape: BaseShape) -> bool:
    return bool(shape.is_placeholder and shape.placeholder_format.type in _TITLE_PLACEHOLDER_TYPES)


def extract_pptx_metadata(content: bytes) -> dict[str, object]:
    """Raw core_properties values, uncleaned -- see extraction_metadata.py
    for why cleaning is centralized there instead of here."""
    props = Presentation(io.BytesIO(content)).core_properties
    return {
        "title": props.title,
        "author": props.author,
        "created_at": props.created.isoformat() if props.created else None,
        "modified_at": props.modified.isoformat() if props.modified else None,
        "language": props.language,
    }


def extract_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    slide_texts: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                markdown_table = rows_to_markdown(
                    [[cell.text for cell in row.cells] for row in shape.table.rows]
                )
                if markdown_table:
                    parts.append(markdown_table)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                parts.append(f"# {text}" if _is_title_shape(shape) else text)
        if parts:
            slide_texts.append("\n\n".join(parts))
    return "\n\n".join(slide_texts)
