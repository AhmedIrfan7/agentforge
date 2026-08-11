"""Unit tests for extraction_pptx.py (roadmap step 092) -- real .pptx
files built with python-pptx itself.
"""

import io
from collections.abc import Callable

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches

from extraction_pptx import extract_pptx, extract_pptx_metadata


def _build_pptx(build: Callable[[PresentationType], None]) -> bytes:
    presentation = Presentation()
    build(presentation)
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def test_title_and_content_layout_title_becomes_h1() -> None:
    def build(p: PresentationType) -> None:
        slide = p.slides.add_slide(p.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text_frame.text = "Body content here."

    text = extract_pptx(_build_pptx(build))
    assert "# Slide Title" in text
    assert "Body content here." in text


def test_title_slide_layout_title_also_becomes_h1() -> None:
    # A dedicated title-slide layout's title is a CENTER_TITLE
    # placeholder, not TITLE -- verified live against a real generated
    # pptx before trusting it; missing this case would silently drop
    # every title-slide-style title.
    def build(p: PresentationType) -> None:
        slide = p.slides.add_slide(p.slide_layouts[0])
        slide.shapes.title.text = "Main Presentation Title"

    text = extract_pptx(_build_pptx(build))
    assert "# Main Presentation Title" in text


def test_table_becomes_markdown_pipe_table_without_duplication() -> None:
    def build(p: PresentationType) -> None:
        slide = p.slides.add_slide(p.slide_layouts[5])
        slide.shapes.title.text = "Scores"
        table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(4), Inches(1))
        rows = [("Name", "Score"), ("Alice", "95"), ("Bob", "80")]
        for row_cells, (name, score) in zip(table_shape.table.rows, rows, strict=True):
            row_cells.cells[0].text = name
            row_cells.cells[1].text = score

    text = extract_pptx(_build_pptx(build))
    assert "| Name | Score |" in text
    assert "| Alice | 95 |" in text
    assert text.count("Alice") == 1


def test_multiple_slides_both_appear_in_order() -> None:
    def build(p: PresentationType) -> None:
        slide1 = p.slides.add_slide(p.slide_layouts[5])
        slide1.shapes.title.text = "First Slide"
        slide2 = p.slides.add_slide(p.slide_layouts[5])
        slide2.shapes.title.text = "Second Slide"

    text = extract_pptx(_build_pptx(build))
    assert text.index("First Slide") < text.index("Second Slide")


def test_empty_presentation_does_not_crash() -> None:
    def build(p: PresentationType) -> None:
        pass

    content = _build_pptx(build)
    assert isinstance(extract_pptx(content), str)


def test_metadata_reads_real_core_properties() -> None:
    def build(p: PresentationType) -> None:
        p.core_properties.title = "Explicit Deck Title"
        p.core_properties.author = "Real Presenter"
        p.core_properties.language = "fr-FR"

    metadata = extract_pptx_metadata(_build_pptx(build))
    assert metadata["title"] == "Explicit Deck Title"
    assert metadata["author"] == "Real Presenter"
    assert metadata["language"] == "fr-FR"
