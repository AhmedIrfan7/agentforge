"""Tests for extraction.py (roadmap step 090) -- exercises
_run_extraction directly against real Postgres + real MinIO (same "no
mocks for infrastructure this project owns" reasoning as
test_document_endpoints.py), not through a real Celery worker: nothing
here needs to prove Celery's own broker/worker plumbing works (step 089
already did, both in test_celery_app.py and live verification) -- only
that this task's own logic (routing by extension, status transitions,
storing extracted text) is correct. The real .delay() round trip through
routers/document.py:upload_document is covered by live verification
instead, same split step 089 used.
"""

import hashlib
import io
import uuid

import docx
import openpyxl
import pytest
from pptx import Presentation
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate
from sqlalchemy import select

from db import get_session, set_tenant_context
from extraction import HANDLERS, _DocumentNotFoundYet, _run_extraction
from models.document import Document
from models.knowledge_base import KnowledgeBase
from models.organization import Organization
from models.workspace import Workspace
from storage import _client_context, ensure_bucket_exists, upload_file


async def _new_org_workspace_kb(slug: str) -> tuple[uuid.UUID, uuid.UUID]:
    async with get_session() as session:
        org = Organization(name="Extraction Test Org", slug=f"{slug}-org")
        session.add(org)
        await session.flush()
        await set_tenant_context(session, org.id)

        workspace = Workspace(tenant_id=org.id, name="Extraction WS", slug=f"{slug}-ws")
        session.add(workspace)
        await session.flush()

        knowledge_base = KnowledgeBase(
            tenant_id=org.id, workspace_id=workspace.id, name="Extraction KB", slug=f"{slug}-kb"
        )
        session.add(knowledge_base)
        await session.flush()
        await session.commit()
        return org.id, knowledge_base.id


async def _create_document(
    tenant_id: uuid.UUID, knowledge_base_id: uuid.UUID, *, title: str, content: bytes
) -> tuple[uuid.UUID, str]:
    storage_key = f"{tenant_id}/{knowledge_base_id}/{uuid.uuid4()}/{title}"
    await ensure_bucket_exists()
    await upload_file(key=storage_key, content=content, content_type="text/plain")

    async with get_session() as session:
        await set_tenant_context(session, tenant_id)
        document = Document(
            tenant_id=tenant_id,
            knowledge_base_id=knowledge_base_id,
            title=title,
            storage_key=storage_key,
            content_type="text/plain",
            size_bytes=len(content),
        )
        session.add(document)
        await session.commit()
        return document.id, storage_key


async def _cleanup(tenant_id: uuid.UUID, storage_key: str | None = None) -> None:
    if storage_key is not None:
        async with _client_context() as s3:
            await s3.delete_object(Bucket="agentforge-dev", Key=storage_key)
    async with get_session() as session:
        await set_tenant_context(session, tenant_id)
        for model in (Document, KnowledgeBase, Workspace):
            result = await session.execute(select(model).where(model.tenant_id == tenant_id))
            for row in result.scalars().all():
                await session.delete(row)
            await session.flush()
        org = await session.get(Organization, tenant_id)
        if org is not None:
            await session.delete(org)
        await session.commit()


def test_expected_extensions_are_registered() -> None:
    # Every extension validation.py:ALLOWED_EXTENSIONS accepts now has a
    # real handler (step 093) -- see extraction.py's own module
    # docstring for why md ended up grouped with the plain-text
    # extensions rather than getting its own transformation.
    assert set(HANDLERS.keys()) == {
        "csv",
        "txt",
        "json",
        "xml",
        "md",
        "pdf",
        "docx",
        "pptx",
        "xlsx",
        "html",
    }


@pytest.mark.anyio
async def test_supported_type_extracts_and_marks_extracted() -> None:
    tenant_id, kb_id = await _new_org_workspace_kb("extract-supported")
    storage_key = None
    try:
        content = b"line one\nline two\n"
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="notes.txt", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            document = await session.get(Document, document_id)
            assert document is not None
            assert document.status == "extracted"
            assert document.extracted_text == content.decode("utf-8")
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_plain_text_document_gets_language_detected_through_the_real_dispatcher() -> None:
    """Plain-text extensions have no format-level metadata source (see
    extraction_metadata.py's module docstring) -- language detection
    from the extracted text is the only doc_metadata field that ever
    populates for them, proven here end-to-end rather than assumed from
    test_extraction_metadata.py's direct unit coverage of detect_language
    alone."""
    tenant_id, kb_id = await _new_org_workspace_kb("extract-language")
    storage_key = None
    try:
        content = (
            b"This is a genuinely ordinary English paragraph about quarterly "
            b"revenue, written with enough real words for confident language "
            b"detection to work reliably."
        )
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="report.txt", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            document = await session.get(Document, document_id)
            assert document is not None
            assert document.status == "extracted"
            assert document.doc_metadata["language"] == "en"
            assert document.doc_metadata["title"] is None
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_unsupported_type_marks_extraction_unsupported() -> None:
    # Every extension validation.py:ALLOWED_EXTENSIONS accepts has a
    # real handler as of step 093 -- "extraction_unsupported" is no
    # longer reachable through any upload the API actually allows, but
    # _run_extraction doesn't re-check ALLOWED_EXTENSIONS itself, only
    # HANDLERS, so this constructs a document with an extension no real
    # upload could ever carry to exercise that fallback path directly.
    tenant_id, kb_id = await _new_org_workspace_kb("extract-unsupported")
    storage_key = None
    try:
        content = b"whatever content, never reaches a handler"
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="mystery.notarealext", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            document = await session.get(Document, document_id)
            assert document is not None
            assert document.status == "extraction_unsupported"
            assert document.extracted_text is None
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_missing_document_raises_not_found_yet() -> None:
    tenant_id, _kb_id = await _new_org_workspace_kb("extract-missing")
    try:
        with pytest.raises(_DocumentNotFoundYet):
            await _run_extraction(uuid.uuid4(), tenant_id)
    finally:
        await _cleanup(tenant_id)


@pytest.mark.anyio
async def test_pdf_document_extracts_through_the_real_dispatcher() -> None:
    """extraction_pdf.py's own logic is unit-tested directly in
    test_extraction_pdf.py -- this only needs to prove HANDLERS["pdf"] is
    actually wired to it end-to-end through _run_extraction, real
    Postgres + real MinIO included."""
    buf = io.BytesIO()
    styles = getSampleStyleSheet()
    SimpleDocTemplate(buf).build(
        [
            Paragraph("Wired Up", styles["Title"]),
            Paragraph("This came from a real PDF through the real dispatcher.", styles["Normal"]),
        ]
    )
    pdf_bytes = buf.getvalue()

    tenant_id, kb_id = await _new_org_workspace_kb("extract-pdf")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.pdf", content=pdf_bytes
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            document = await session.get(Document, document_id)
            assert document is not None
            assert document.status == "extracted"
            assert document.extracted_text is not None
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_docx_document_extracts_through_the_real_dispatcher() -> None:
    """extraction_docx.py's own logic is unit-tested directly in
    test_extraction_docx.py -- this only needs to prove HANDLERS["docx"]
    is actually wired to it end-to-end through _run_extraction, and
    (roadmap step 094) that doc_metadata gets populated in the same
    pass -- covered here rather than in every other format's wiring
    test, since extraction_metadata.py's own merge/cleaning logic is
    already covered directly in test_extraction_metadata.py."""
    document = docx.Document()
    document.core_properties.title = "Wired Doc Title"
    document.core_properties.author = "Wired Author"
    document.add_paragraph("This came from a real DOCX through the real dispatcher.")
    buf = io.BytesIO()
    document.save(buf)

    tenant_id, kb_id = await _new_org_workspace_kb("extract-docx")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.docx", content=buf.getvalue()
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.extracted_text == (
                "This came from a real DOCX through the real dispatcher."
            )
            assert fetched.doc_metadata["title"] == "Wired Doc Title"
            assert fetched.doc_metadata["author"] == "Wired Author"
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_pptx_document_extracts_through_the_real_dispatcher() -> None:
    """extraction_pptx.py's own logic is unit-tested directly in
    test_extraction_pptx.py -- this only needs to prove HANDLERS["pptx"]
    is actually wired to it end-to-end through _run_extraction."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Wired Up"
    buf = io.BytesIO()
    presentation.save(buf)

    tenant_id, kb_id = await _new_org_workspace_kb("extract-pptx")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.pptx", content=buf.getvalue()
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.extracted_text == "# Wired Up"
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_xlsx_document_extracts_through_the_real_dispatcher() -> None:
    """extraction_xlsx.py's own logic is unit-tested directly in
    test_extraction_xlsx.py -- this only needs to prove HANDLERS["xlsx"]
    is actually wired to it end-to-end through _run_extraction."""
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Data"
    sheet.append(["Name", "Score"])
    sheet.append(["Alice", 95])
    buf = io.BytesIO()
    workbook.save(buf)

    tenant_id, kb_id = await _new_org_workspace_kb("extract-xlsx")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.xlsx", content=buf.getvalue()
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.extracted_text is not None
            assert "# Data" in fetched.extracted_text
            assert "| Alice | 95 |" in fetched.extracted_text
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_html_document_extracts_through_the_real_dispatcher() -> None:
    """extraction_html.py's own logic is unit-tested directly in
    test_extraction_html.py -- this only needs to prove HANDLERS["html"]
    is actually wired to it end-to-end through _run_extraction."""
    content = b"<h1>Wired Up</h1><p>This came from real HTML through the real dispatcher.</p>"

    tenant_id, kb_id = await _new_org_workspace_kb("extract-html")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.html", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.extracted_text is not None
            assert "# Wired Up" in fetched.extracted_text
            assert "This came from real HTML through the real dispatcher." in (
                fetched.extracted_text
            )
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_md_document_extracts_through_the_real_dispatcher_as_is() -> None:
    """md is registered to the same plain-text passthrough as csv/txt/
    json/xml (see extraction.py's module docstring for why) -- this
    proves HANDLERS["md"] is wired to it, and that an uploaded markdown
    file's bytes come back completely unchanged, not reformatted."""
    content = b"# Already Markdown\n\nThis file is already in the target format."

    tenant_id, kb_id = await _new_org_workspace_kb("extract-md")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.md", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.extracted_text == content.decode("utf-8")
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_document_type_is_populated_through_the_real_dispatcher() -> None:
    """agents/document_analysis.py's own classification logic is
    unit-tested directly in test_document_analysis_agent.py -- this only
    needs to prove _run_extraction actually calls it and stores the
    result in doc_metadata."""
    content = b"# Frequently Asked Questions\n\nQ: How do I reset my password?\nA: Click the link."

    tenant_id, kb_id = await _new_org_workspace_kb("extract-doctype")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.txt", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.doc_metadata["document_type"] == "faq"
            signals = fetched.doc_metadata["document_type_signals"]
            assert isinstance(signals, list)
            assert "frequently asked questions" in signals
    finally:
        await _cleanup(tenant_id, storage_key)


@pytest.mark.anyio
async def test_quality_signals_are_populated_through_the_real_dispatcher() -> None:
    """quality.py's own checks are unit-tested directly in
    test_quality.py -- this only needs to prove _run_extraction actually
    calls assess_quality and stores content_hash + the quality flags."""
    content = b"Real, clean content for the quality-signal wiring test."

    tenant_id, kb_id = await _new_org_workspace_kb("extract-quality")
    storage_key = None
    try:
        document_id, storage_key = await _create_document(
            tenant_id, kb_id, title="wired.txt", content=content
        )

        await _run_extraction(document_id, tenant_id)

        async with get_session() as session:
            await set_tenant_context(session, tenant_id)
            fetched = await session.get(Document, document_id)
            assert fetched is not None
            assert fetched.status == "extracted"
            assert fetched.content_hash == hashlib.sha256(content).hexdigest()
            quality = fetched.doc_metadata["quality"]
            assert quality == {"is_empty": False, "has_broken_formatting": False}
    finally:
        await _cleanup(tenant_id, storage_key)
